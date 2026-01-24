"""
é€šç”¨æ–‡æ¡£åŠ è½½å™¨

æ”¯æŒ Markdownã€TXTã€PPTXã€PDFã€DOCXã€DOC ç­‰å¤šç§æ ¼å¼ã€‚
ç¬¦åˆ LangChain æ–‡æ¡£åŠ è½½å™¨æ¥å£è§„èŒƒï¼Œæ”¯æŒæŒ‰ç±»åˆ«ï¼ˆpolicies/casesï¼‰ç»„ç»‡çŸ¥è¯†åº“ã€‚
æ‰€æœ‰æ ¼å¼éƒ½ä¼šç»Ÿä¸€è½¬æ¢ä¸º Markdownï¼Œå¹¶è‡ªåŠ¨æ¸…ç†å†—ä½™ä¿¡æ¯ã€‚
"""
import re
import subprocess
from pathlib import Path
from typing import Optional, Literal

from langchain_core.documents import Document
import filetype

# Optional imports for specific file formats
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None


# ==================== æ–‡ä»¶ç±»å‹æ£€æµ‹ ====================

class FileTypeDetector:
    """é€šè¿‡æ–‡ä»¶å†…å®¹æ£€æµ‹çœŸå®ç±»å‹ï¼Œä¸ä¾èµ–æ‰©å±•å"""

    MIME_TYPE_MAP = {
        'application/pdf': 'pdf',
        'application/msword': 'doc',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/vnd.ms-word': 'doc',
        'application/wps-office.doc': 'doc',
        'application/vnd.ms-powerpoint': 'ppt',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'application/vnd.ms-excel': 'xls',
        'application/wps-office.xls': 'xls',
        'text/plain': 'txt',
        'text/markdown': 'md',
        'application/x-ole-storage': 'ole',
    }

    EXT_TYPE_MAP = {
        '.pdf': 'pdf', '.doc': 'doc', '.docx': 'docx',
        '.ppt': 'ppt', '.pptx': 'pptx', '.txt': 'txt',
        '.md': 'markdown',
    }

    @classmethod
    def detect(cls, file_path: Path) -> str:
        """æ£€æµ‹æ–‡ä»¶çœŸå®ç±»å‹"""
        ext = file_path.suffix.lower()
        kind = filetype.guess(str(file_path))

        if kind is None:
            return cls.EXT_TYPE_MAP.get(ext, 'unknown')

        mime_type = kind.mime
        doc_type = cls.MIME_TYPE_MAP.get(mime_type)

        if doc_type == 'ole':
            return 'doc' if ext in ['.doc', '.docx'] else doc_type

        if ext == '.docx' and mime_type in [
            'application/vnd.ms-excel',
            'application/wps-office.xls',
            'application/x-ole-storage'
        ]:
            print(f"âš ï¸  æ£€æµ‹åˆ°ä¼ªè£…æˆ .docx çš„ .doc æ–‡ä»¶: {file_path.name}")
            return 'doc'

        return doc_type or cls.EXT_TYPE_MAP.get(ext, 'unknown')


# ==================== Markdown æ¸…ç† ====================

class MarkdownCleaner:
    """æ¸…ç† Markdown å†…å®¹ï¼Œå»é™¤æ ¼å¼æ•°æ®å’Œå†—ä½™ä¿¡æ¯"""

    FOOTER_PATTERNS = [
        r'ç¬¬\s*\d+\s*é¡µ', r'Page\s*\d+',
        r'ä¿å¯†|æœºå¯†|å†…éƒ¨èµ„æ–™', r'www\.\w+\.com', r'http[s]?://\S+',
    ]

    PLACEHOLDER_PATTERNS = [
        r'ç‚¹å‡»æ­¤å¤„æ·»åŠ .*', r'è¯·è¾“å…¥.*', r'\[.*?\]', r'{{.*?}}',
    ]

    @classmethod
    def clean_text(cls, text: str) -> str:
        """æ¸…ç†æ–‡æœ¬å†…å®¹"""
        for pattern in cls.FOOTER_PATTERNS + cls.PLACEHOLDER_PATTERNS:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE)

        text = re.sub(r'\n\s*\n\s*\n+', '\n\n\n', text)
        lines = [line.strip() for line in text.split('\n')]

        cleaned_lines = []
        for line in lines:
            chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', line))
            alnum_chars = len(re.findall(r'[a-zA-Z0-9]', line))
            total_chars = len(line)

            if (chinese_chars + alnum_chars) / max(total_chars, 1) > 0.1 or (chinese_chars + alnum_chars) >= 5:
                cleaned_lines.append(line)

        return '\n'.join(cleaned_lines).strip()

    @classmethod
    def is_meaningful_content(cls, text: str, min_length: int = 20) -> bool:
        """åˆ¤æ–­æ–‡æœ¬æ˜¯å¦æœ‰æ„ä¹‰"""
        if len(text.strip()) < min_length:
            return False
        if not re.search(r'[\u4e00-\u9fff\u4e00-\u9fa5a-zA-Z0-9]', text):
            return False
        if text.strip().replace('-', '').replace('/', '').replace(':', '').strip().isdigit():
            return False
        return True


# ==================== åŸºç¡€åŠ è½½å™¨ ====================

class BaseDocumentLoader:
    """æ–‡æ¡£åŠ è½½å™¨åŸºç±»ï¼Œæä¾›å…¬å…±æ–¹æ³•"""

    def __init__(
        self,
        file_path: str | Path,
        category: Optional[Literal["policies", "cases"]] = None,
    ):
        self.file_path = Path(file_path)
        self.category = category
        self.cleaner = MarkdownCleaner()

    def _create_document(
        self,
        content: str,
        **metadata_kwargs
    ) -> Optional[Document]:
        """åˆ›å»ºæ–‡æ¡£å¯¹è±¡"""
        if not self.cleaner.is_meaningful_content(content):
            return None

        cleaned_content = self.cleaner.clean_text(content)
        if not self.cleaner.is_meaningful_content(cleaned_content):
            return None

        metadata = {
            "source": str(self.file_path.name),
            **metadata_kwargs
        }

        if self.category:
            metadata["category"] = self.category

        return Document(page_content=cleaned_content, metadata=metadata)

    def _validate_file(self) -> None:
        """éªŒè¯æ–‡ä»¶å­˜åœ¨"""
        if not self.file_path.exists():
            raise FileNotFoundError(f"æ‰¾ä¸åˆ°æ–‡ä»¶: {self.file_path}")


# ==================== DOC åŠ è½½å™¨ ====================

class DOCLoader(BaseDocumentLoader):
    """Word æ–‡æ¡£åŠ è½½å™¨ï¼ˆDOC - Legacy Office 97-2003 æ ¼å¼ï¼‰"""

    def load(self) -> list[Document]:
        """åŠ è½½ DOC æ–‡ä»¶"""
        self._validate_file()
        print(f"ğŸ“ æ­£åœ¨è¯»å– Word æ–‡æ¡£ï¼ˆDOC æ ¼å¼ï¼‰: {self.file_path} ...")

        try:
            text = self._extract_text()
            return self._parse_text(text)
        except Exception as e:
            raise Exception(f"è¯»å– Word æ–‡æ¡£ï¼ˆDOCï¼‰å¤±è´¥: {e}\næç¤º: è¯·å®‰è£… catdoc: sudo apt-get install catdoc")

    def _extract_text(self) -> str:
        """å°è¯•å¤šç§æ–¹æ³•æå–æ–‡æœ¬"""
        extractors = [
            self._extract_with_antiword,
            self._extract_with_catdoc,
            self._extract_with_olefile,
        ]

        for idx, extractor in enumerate(extractors, 1):
            try:
                if idx > 1:
                    print(f"âš ï¸  å‰ä¸€ç§æ–¹æ³•å¤±è´¥ï¼Œå°è¯•ä¸‹ä¸€ç§...")
                return extractor()
            except Exception as e:
                continue

        raise Exception("æ‰€æœ‰æå–æ–¹æ³•éƒ½å¤±è´¥")

    def _extract_with_antiword(self) -> str:
        """ä½¿ç”¨ antiword æå–æ–‡æœ¬"""
        try:
            result = subprocess.run(
                ['antiword', str(self.file_path)],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                return result.stdout
            raise Exception(f"antiword failed with code {result.returncode}")
        except FileNotFoundError:
            raise Exception("antiword æœªå®‰è£…ï¼Œè¯·è¿è¡Œ: sudo apt-get install antiword")

    def _extract_with_catdoc(self) -> str:
        """ä½¿ç”¨ catdoc æå–æ–‡æœ¬"""
        try:
            result = subprocess.run(
                ['catdoc', str(self.file_path)],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                return result.stdout
            raise Exception(f"catdoc failed with code {result.returncode}")
        except FileNotFoundError:
            raise Exception("catdoc æœªå®‰è£…ï¼Œè¯·è¿è¡Œ: sudo apt-get install catdoc")

    def _extract_with_olefile(self) -> str:
        """ä½¿ç”¨ olefile æå–æ–‡æœ¬ï¼ˆçº¯ Python å¤‡ç”¨æ–¹æ¡ˆï¼‰"""
        try:
            import olefile
        except ImportError:
            raise Exception("olefile æœªå®‰è£…ï¼Œè¯·è¿è¡Œ: pip install olefile")

        ole = olefile.OleFileIO(self.file_path)

        if not ole.exists('WordDocument'):
            raise Exception("ä¸æ˜¯æœ‰æ•ˆçš„ Word æ–‡æ¡£")

        text_parts = []
        streams_to_try = ['WordDocument', '1Table', '0Table', 'Data']

        for stream_name in streams_to_try:
            if not ole.exists(stream_name):
                continue

            try:
                data = ole.openstream(stream_name).read()
                print(f"   æ­£åœ¨è§£ææµ: {stream_name} ({len(data)} å­—èŠ‚)")

                for encoding in ['utf-16le', 'utf-16be', 'gbk', 'gb2312', 'gb18030', 'utf-8']:
                    try:
                        decoded = data.decode(encoding, errors='ignore')
                        chinese_chars = re.findall(r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+', decoded)
                        if chinese_chars:
                            text_segment = ''.join(chinese_chars)
                            text_segment = re.sub(r'\s+', '\n', text_segment)
                            text_segment = re.sub(r'\n{3,}', '\n\n', text_segment)

                            if len(text_segment.strip()) > 20:
                                print(f"   âœ“ ä½¿ç”¨ {encoding} ç¼–ç æå–äº† {len(text_segment)} ä¸ªå­—ç¬¦")
                                text_parts.append(text_segment)
                                break
                    except:
                        continue
            except Exception as e:
                print(f"   âš ï¸  è§£ææµ {stream_name} å¤±è´¥: {e}")
                continue

        ole.close()

        if not text_parts:
            raise Exception("æœªèƒ½ä»æ–‡æ¡£ä¸­æå–åˆ°æœ‰æ•ˆæ–‡æœ¬ï¼Œè¯·å°†æ–‡ä»¶è½¬æ¢ä¸º .docx æ ¼å¼")

        return '\n\n'.join(text_parts)

    def _parse_text(self, text: str) -> list[Document]:
        """è§£ææå–çš„æ–‡æœ¬"""
        cleaned_text = self.cleaner.clean_text(text)
        paragraphs = [p.strip() for p in cleaned_text.split('\n\n') if p.strip()]

        documents = []
        for idx, paragraph in enumerate(paragraphs, start=1):
            doc = self._create_document(
                f"## æ®µè½ {idx}\n\n{paragraph}",
                paragraph=idx,
                type="doc",
            )
            if doc:
                documents.append(doc)

        print(f"âœ… æå–å®Œæˆï¼Œå…±è·å– {len(documents)} ä¸ªæ®µè½")
        return documents


# ==================== DOCX åŠ è½½å™¨ ====================

class DOCXLoader(BaseDocumentLoader):
    """Word æ–‡æ¡£åŠ è½½å™¨ï¼ˆDOCXï¼‰"""

    def load(self) -> list[Document]:
        """åŠ è½½ DOCX æ–‡ä»¶"""
        self._validate_file()
        print(f"ğŸ“ æ­£åœ¨è¯»å– Word æ–‡æ¡£ï¼ˆDOCX æ ¼å¼ï¼‰: {self.file_path} ...")

        if DocxDocument is None:
            raise Exception("python-docx æœªå®‰è£…ï¼Œè¯·è¿è¡Œ: pip install python-docx")

        try:
            doc = DocxDocument(str(self.file_path))
            return self._parse_docx(doc)
        except Exception as e:
            print(f"âš ï¸  python-docx è¯»å–å¤±è´¥: {e}")
            print(f"âš ï¸  å¯èƒ½æ˜¯ä¼ªè£…æˆ .docx çš„ .doc æ–‡ä»¶ï¼Œå°è¯•ä½¿ç”¨ DOCLoader...")
            return DOCLoader(self.file_path, self.category).load()

    def _parse_docx(self, doc) -> list[Document]:
        """è§£æ DOCX æ–‡æ¡£"""
        documents = []
        current_content = []
        current_heading = "æ–‡æ¡£å¼€å§‹"
        current_level = 0
        paragraph_count = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else ""

            if 'Heading' in style_name:
                if current_content:
                    doc = self._create_document(
                        '\n'.join(current_content).strip(),
                        paragraph=paragraph_count,
                        type="docx",
                    )
                    if doc:
                        documents.append(doc)

                current_content = []
                current_heading = text
                current_level = self._extract_heading_level(style_name)
                paragraph_count += 1
            else:
                current_content.append(text)

        if current_content:
            doc = self._create_document(
                '\n'.join(current_content).strip(),
                paragraph=paragraph_count,
                type="docx",
            )
            if doc:
                documents.append(doc)

        if not documents:
            return self._parse_as_paragraphs(doc)

        print(f"âœ… æå–å®Œæˆï¼Œå…±è·å– {len(documents)} ä¸ªæ–‡æ¡£ç‰‡æ®µ")
        return documents

    def _extract_heading_level(self, style_name: str) -> int:
        """ä»æ ·å¼åç§°æå–æ ‡é¢˜çº§åˆ«"""
        for i in range(1, 7):
            if f'Heading {i}' in style_name:
                return i
        return 6

    def _parse_as_paragraphs(self, doc) -> list[Document]:
        """æŒ‰æ®µè½è§£æï¼ˆå½“æ²¡æœ‰æ£€æµ‹åˆ°æ ‡é¢˜æ—¶ï¼‰"""
        documents = []
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        for idx, para_text in enumerate(paragraphs, start=1):
            doc = self._create_document(
                f"## æ®µè½ {idx}\n\n{para_text}",
                paragraph=idx,
                type="docx",
            )
            if doc:
                documents.append(doc)

        return documents


# ==================== PDF åŠ è½½å™¨ ====================

class PDFLoader(BaseDocumentLoader):
    """PDF æ–‡æ¡£åŠ è½½å™¨"""

    def load(self) -> list[Document]:
        """åŠ è½½ PDF æ–‡ä»¶"""
        self._validate_file()
        print(f"ğŸ“„ æ­£åœ¨è¯»å– PDF: {self.file_path} ...")

        if PdfReader is None:
            raise Exception("pypdf æœªå®‰è£…ï¼Œè¯·è¿è¡Œ: pip install pypdf")

        try:
            reader = PdfReader(str(self.file_path))
            return self._parse_pdf(reader)
        except Exception as e:
            raise Exception(f"è¯»å– PDF æ–‡ä»¶å¤±è´¥: {e}")

    def _parse_pdf(self, reader) -> list[Document]:
        """è§£æ PDF æ–‡æ¡£"""
        documents = []

        for page_idx, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text()
                if not text or not text.strip():
                    continue

                doc = self._create_document(
                    f"# ç¬¬ {page_idx} é¡µ\n\n{text}",
                    page=page_idx,
                    type="pdf",
                )
                if doc:
                    documents.append(doc)
            except Exception as e:
                print(f"âš ï¸  å¤„ç†ç¬¬ {page_idx} é¡µæ—¶å‡ºé”™: {e}")
                continue

        print(f"âœ… æå–å®Œæˆï¼Œå…±è·å– {len(documents)} é¡µæœ‰æ•ˆå†…å®¹")
        return documents


# ==================== PPTX åŠ è½½å™¨ ====================

class PPTXLoader(BaseDocumentLoader):
    """PPTX æ–‡æ¡£åŠ è½½å™¨"""

    def load(self) -> list[Document]:
        """åŠ è½½ PPTX æ–‡ä»¶"""
        self._validate_file()
        print(f"ğŸ“‚ æ­£åœ¨è¯»å– PPT: {self.file_path} ...")

        if Presentation is None:
            raise Exception("python-pptx æœªå®‰è£…ï¼Œè¯·è¿è¡Œ: pip install python-pptx")

        prs = Presentation(str(self.file_path))
        documents = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]

            if slide_texts:
                content = "\n".join(slide_texts)
                doc = self._create_document(
                    f"# ç¬¬ {slide_idx} é¡µ\n\n{content}",
                    page=slide_idx,
                    type="pptx",
                )
                if doc:
                    documents.append(doc)

        print(f"âœ… æå–å®Œæˆï¼Œå…±è·å– {len(documents)} é¡µæœ‰æ•ˆå†…å®¹")
        return documents


# ==================== Markdown åŠ è½½å™¨ ====================

class MarkdownLoader(BaseDocumentLoader):
    """Markdown æ–‡æ¡£åŠ è½½å™¨"""

    def __init__(
        self,
        file_path: str | Path,
        encoding: str = "utf-8",
        category: Optional[Literal["policies", "cases"]] = None,
    ):
        super().__init__(file_path, category)
        self.encoding = encoding

    def load(self) -> list[Document]:
        """åŠ è½½ Markdown æ–‡ä»¶"""
        self._validate_file()
        print(f"ğŸ“„ æ­£åœ¨è¯»å– Markdown: {self.file_path} ...")

        with open(self.file_path, "r", encoding=self.encoding, errors="ignore") as f:
            content = f.read()

        content = self.cleaner.clean_text(content)
        documents = self._split_by_headers(content)

        print(f"âœ… æå–å®Œæˆï¼Œå…±è·å– {len(documents)} ä¸ªæ–‡æ¡£ç‰‡æ®µ")
        return documents

    def _split_by_headers(self, content: str) -> list[Document]:
        """æŒ‰ Markdown æ ‡é¢˜åˆ†å‰²æ–‡æ¡£"""
        documents = []
        lines = content.split("\n")

        current_section = []
        current_header = "æ–‡æ¡£å¼€å§‹"
        current_level = 0
        section_idx = 0

        for line in lines:
            header_match = re.match(r'^(#{1,6})\s+(.+)$', line)

            if header_match:
                if current_section:
                    section_content = "\n".join(current_section).strip()
                    doc = self._create_document(
                        section_content,
                        section=section_idx + 1,
                        type="markdown",
                        header=current_header,
                        header_level=current_level,
                    )
                    if doc:
                        documents.append(doc)
                        section_idx += 1

                current_level = len(header_match.group(1))
                current_header = header_match.group(2).strip()
                current_section = [line]
            else:
                current_section.append(line)

        if current_section:
            section_content = "\n".join(current_section).strip()
            doc = self._create_document(
                section_content,
                section=section_idx + 1,
                type="markdown",
                header=current_header,
                header_level=current_level,
            )
            if doc:
                documents.append(doc)

        return documents


# ==================== æ–‡æœ¬æ–‡ä»¶åŠ è½½å™¨ ====================

class TextFileLoader(BaseDocumentLoader):
    """TXT æ–‡æ¡£åŠ è½½å™¨"""

    def __init__(
        self,
        file_path: str | Path,
        encoding: str = "utf-8",
        category: Optional[Literal["policies", "cases"]] = None,
    ):
        super().__init__(file_path, category)
        self.encoding = encoding

    def load(self) -> list[Document]:
        """åŠ è½½æ–‡æœ¬æ–‡ä»¶"""
        self._validate_file()
        print(f"ğŸ“‚ æ­£åœ¨è¯»å–æ–‡æœ¬æ–‡ä»¶: {self.file_path} ...")

        with open(self.file_path, "r", encoding=self.encoding, errors="ignore") as f:
            content = f.read()

        content = self.cleaner.clean_text(content)
        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]

        documents = []
        for idx, paragraph in enumerate(paragraphs, start=1):
            doc = self._create_document(
                f"## æ®µè½ {idx}\n\n{paragraph}",
                paragraph=idx,
                type="text",
            )
            if doc:
                documents.append(doc)

        print(f"âœ… æå–å®Œæˆï¼Œå…±è·å– {len(documents)} ä¸ªæ®µè½")
        return documents


# ==================== æ‰¹é‡åŠ è½½å‡½æ•° ====================

def load_documents_from_directory(
    directory: str | Path,
    file_extensions: Optional[list[str]] = None,
    category: Optional[Literal["policies", "cases"]] = None,
) -> list[Document]:
    """ä»ç›®å½•æ‰¹é‡åŠ è½½æ–‡æ¡£ï¼Œè‡ªåŠ¨æ£€æµ‹çœŸå®æ–‡ä»¶ç±»å‹"""
    directory = Path(directory)
    if not directory.exists():
        raise FileNotFoundError(f"ç›®å½•ä¸å­˜åœ¨: {directory}")

    if file_extensions is None:
        file_extensions = [".md", ".txt", ".pptx", ".ppt", ".pdf", ".docx", ".doc"]

    all_documents = []

    for file_path in directory.rglob("*"):
        if not file_path.is_file() or file_path.suffix.lower() not in file_extensions:
            continue

        try:
            real_type = FileTypeDetector.detect(file_path)
            print(f"ğŸ” æ£€æµ‹æ–‡ä»¶ç±»å‹: {file_path.name} -> {real_type}")

            loader = _create_loader(file_path, real_type, category)
            if loader is None:
                continue

            documents = loader.load()
            all_documents.extend(documents)
        except Exception as e:
            print(f"âš ï¸  åŠ è½½æ–‡ä»¶ {file_path} æ—¶å‡ºé”™: {e}")
            continue

    print(f"ğŸ“š æ€»å…±åŠ è½½äº† {len(all_documents)} ä¸ªæ–‡æ¡£ç‰‡æ®µ")
    return all_documents


def _create_loader(
    file_path: Path,
    file_type: str,
    category: Optional[Literal["policies", "cases"]],
) -> Optional[BaseDocumentLoader]:
    """æ ¹æ®æ–‡ä»¶ç±»å‹åˆ›å»ºå¯¹åº”çš„åŠ è½½å™¨"""
    loader_map = {
        'pdf': PDFLoader,
        'doc': DOCLoader,
        'docx': DOCXLoader,
        'pptx': PPTXLoader,
        'markdown': MarkdownLoader,
        'txt': TextFileLoader,
    }

    loader_class = loader_map.get(file_type)
    if loader_class is None:
        if file_type == 'ppt':
            print(f"âš ï¸  æš‚ä¸æ”¯æŒ PPT æ ¼å¼ï¼Œè¯·è½¬æ¢ä¸º PPTX: {file_path.name}")
        else:
            print(f"âš ï¸  ä¸æ”¯æŒçš„æ–‡ä»¶ç±»å‹: {file_type}ï¼Œè·³è¿‡: {file_path.name}")
        return None

    return loader_class(file_path, category=category)


def load_knowledge_base(
    data_dir: str | Path,
    categories: Optional[list[Literal["policies", "cases"]]] = None,
) -> list[Document]:
    """åŠ è½½çŸ¥è¯†åº“ï¼ˆæ”¯æŒåˆ†ç±»ï¼‰"""
    data_dir = Path(data_dir)

    if categories is None:
        categories = [
            item.name
            for item in data_dir.iterdir()
            if item.is_dir() and item.name in ["policies", "cases"]
        ]

    if not categories:
        raise FileNotFoundError(
            f"æœªæ‰¾åˆ°ä»»ä½•ç±»åˆ«ç›®å½•ã€‚è¯·åœ¨ {data_dir} ä¸‹åˆ›å»º 'policies' å’Œ/æˆ– 'cases' ç›®å½•ã€‚"
        )

    all_documents = []

    for category in categories:
        category_dir = data_dir / category
        if not category_dir.exists():
            print(f"âš ï¸  ç›®å½•ä¸å­˜åœ¨ï¼Œè·³è¿‡: {category_dir}")
            continue

        print(f"\n{'='*60}")
        print(f"æ­£åœ¨åŠ è½½ç±»åˆ«: {category}")
        print(f"{'='*60}")

        documents = load_documents_from_directory(
            category_dir,
            file_extensions=[".md", ".txt", ".pptx", ".ppt", ".pdf", ".docx", ".doc"],
            category=category,
        )
        all_documents.extend(documents)

    print(f"\n{'='*60}")
    print(f"âœ… çŸ¥è¯†åº“åŠ è½½å®Œæˆï¼")
    print(f"   - æ€»æ–‡æ¡£æ•°: {len(all_documents)}")
    print(f"   - ç±»åˆ«: {', '.join(categories)}")
    print(f"{'='*60}\n")

    return all_documents
