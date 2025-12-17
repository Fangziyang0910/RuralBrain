import os
from pptx import Presentation
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

# ================= 配置区 =================
# 1. 你的 PPT 原始路径 (WSL 下访问 Windows 盘符)
# 1. 获取当前脚本目录
CURRENT_DIR = os.path.dirname(os.path.abspath(__file__))

# 2. 指向刚才复制进来的 PPT (注意文件名变了)
PPT_PATH = os.path.join(CURRENT_DIR, "..", "data", "luofu_strategy.pptx")

# 3. 数据库路径
PROJECT_ROOT = os.path.dirname(os.path.dirname(CURRENT_DIR))
PERSIST_DIRECTORY = os.path.join(PROJECT_ROOT, "knowledge_base", "luofu_db")
# =========================================

def extract_text_from_pptx(file_path):
    """直接从 PPTX 提取文字，不转 PDF"""
    if not os.path.exists(file_path):
        print(f"❌ 错误：找不到文件 {file_path}")
        return []

    print(f"📂 正在读取 PPT: {file_path} ...")
    prs = Presentation(file_path)
    documents = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        # 遍历每页的所有文本框
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        page_content = "\n".join(slide_text)
        
        # 只有当这一页有字的时候才保存
        if page_content:
            # 加上页码元数据，方便以后知道是哪一页
            doc = Document(
                page_content=page_content,
                metadata={"source": "博罗古城总体规划说明书", "page": i + 1}
            )
            documents.append(doc)
    
    print(f"✅ 提取完成，共获取 {len(documents)} 页有效内容。")
    return documents

def build_vector_db():
    # 1. 提取文字
    docs = extract_text_from_pptx(PPT_PATH)
    if not docs:
        return

    # 2. 切分文本
    print("✂️ 正在切分文本...")
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,    # 每块的大小
        chunk_overlap=100  # 重叠部分，防止切断句子
    )
    splits = text_splitter.split_documents(docs)
    print(f"✂️ 切分完成，共 {len(splits)} 个片段。")

    # 3. 向量化并存储
    print("🧠 正在向量化 (首次运行需要下载模型，请稍候)...")
    embedding_model = HuggingFaceEmbeddings(model_name="BAAI/bge-small-zh-v1.5")

    print(f"💾 正在写入数据库: {PERSIST_DIRECTORY}")
    vectorstore = Chroma.from_documents(
        documents=splits,
        embedding=embedding_model,
        persist_directory=PERSIST_DIRECTORY
    )
    print("🎉 恭喜！知识库构建成功！")

if __name__ == "__main__":
    build_vector_db()